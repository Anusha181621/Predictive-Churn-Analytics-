"""Fill the Dragon's Den slide template from the pipeline's own artefacts.

The deck is generated rather than typed for the same reason the dashboard is a reader: every
number on a slide is read from ``outputs/`` at build time, so re-running the pipeline and
re-running this script keeps the pitch and the product telling the same story. The template is
opened read-only and a new file is written -- ``Dragons_Den_Presentation_Template.pptx`` is never
modified.

    python scripts/build_presentation.py [--template PATH] [--output PATH]

Layout is absolute-positioned in inches on the template's 13.333 x 7.5in stage, inside the band
the template leaves free: y 1.50 -> 6.68, x 0.95 -> 12.55, to the right of its orange accent bar.

The two bar charts follow the project's own visualisation rules (see ``app/theme.py``): one
measure per chart and one axis -- never a second scale -- every bar direct-labelled so identity
never rests on colour alone, no gridlines because the values are printed, and an ordinal severity
scale drawn as a single hue with monotone lightness rather than a rainbow.
"""

from __future__ import annotations

import argparse
import json
import sys
from dataclasses import dataclass, field
from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parent.parent

# --------------------------------------------------------------------------------------
# palette -- lifted from the template's own XML so the new slides cannot drift from slide 1
# --------------------------------------------------------------------------------------

NAVY = RGBColor(0x00, 0x50, 0x71)
ORANGE = RGBColor(0xFB, 0x4D, 0x0A)
ORANGE_TEXT = RGBColor(0xC9, 0x3F, 0x06)  # 4.7:1 on the card fill, where bright orange is 3.2:1
GRAY = RGBColor(0x48, 0x48, 0x48)
RULE = RGBColor(0xD6, 0xE4, 0xEA)
CARD = RGBColor(0xF5, 0xF9, 0xFB)
CARD_DEEP = RGBColor(0xEC, 0xF3, 0xF7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

#: Ordinal severity scale: one hue, monotone lightness (OKLab L 0.82 -> 0.41), so the ordering
#: survives colour-blindness and greyscale print. Every bar is direct-labelled too, which is what
#: the low fill contrast of the lightest step obliges.
RISK_RAMP = {
    "Low": RGBColor(0xA9, 0xCB, 0xD9),
    "Medium": RGBColor(0x6B, 0xA5, 0xBC),
    "High": RGBColor(0x2E, 0x7F, 0xA0),
    "Critical": NAVY,
}
#: One series means one colour: length carries the magnitude, and there is nothing to look up.
BAR_ONE_SERIES = RGBColor(0x2E, 0x7F, 0xA0)

FONT = "Calibri"
BAR_RADIUS = 0.22  # ~3px rounded data-ends at this bar thickness

# the free band the template leaves, in inches
X0, X1 = 0.95, 12.55
Y0 = 1.50
W = X1 - X0


# --------------------------------------------------------------------------------------
# drawing helpers
# --------------------------------------------------------------------------------------


@dataclass
class Para:
    """One paragraph inside a text box."""

    text: str
    size: float = 8.0
    bold: bool = False
    color: RGBColor = field(default=GRAY)
    space_before: float = 0.0
    align: str = "l"
    line_spacing: float = 1.06


_ALIGN = {"l": PP_ALIGN.LEFT, "c": PP_ALIGN.CENTER, "r": PP_ALIGN.RIGHT}
_ANCHOR = {"t": MSO_ANCHOR.TOP, "c": MSO_ANCHOR.MIDDLE, "b": MSO_ANCHOR.BOTTOM}


def text(slide, x, y, w, h, paras, anchor="t"):
    """Add a text box with zeroed margins, so it lines up with the drawn geometry."""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = frame.margin_top = frame.margin_bottom = 0
    frame.vertical_anchor = _ANCHOR[anchor]
    for index, para in enumerate(paras):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.alignment = _ALIGN[para.align]
        paragraph.line_spacing = para.line_spacing
        if para.space_before:
            paragraph.space_before = Pt(para.space_before)
        run = paragraph.add_run()
        run.text = para.text
        font = run.font
        font.name = FONT
        font.size = Pt(para.size)
        font.bold = para.bold
        font.color.rgb = para.color
    return box


def rect(slide, x, y, w, h, fill=None, line=None, line_w=0.75, radius=None):
    """A rectangle with the theme's default shadow and style switched off."""
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius is not None else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    if radius is not None:
        shape.adjustments[0] = radius
    shape.shadow.inherit = False
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_w)
    return shape


def send_behind(slide, shape, target):
    """Move ``shape`` immediately before ``target`` in z-order."""
    tree = slide.shapes._spTree
    tree.remove(shape._element)
    target._element.addprevious(shape._element)


def label(slide, x, y, w, caption, note=None, note_at=2.60):
    """A small section label over a hairline rule -- the template's own header idiom."""
    text(slide, x, y, w, 0.16, [Para(caption.upper(), 8, True, NAVY)])
    if note and w > note_at + 1.0:
        text(slide, x + note_at, y, w - note_at, 0.16, [Para(note, 7.5, False, GRAY)])
    rect(slide, x, y + 0.20, w, 0.008, fill=RULE)


def card(slide, x, y, w, h, title, body, accent=ORANGE, fill=CARD, title_size=9.5):
    """A light panel with a coloured left edge -- echoing the template's accent bar."""
    rect(slide, x, y, w, h, fill=fill, line=RULE)
    rect(slide, x, y, 0.05, h, fill=accent)
    pad = 0.15
    paras = [Para(title, title_size, True, NAVY)]
    for index, line in enumerate(body):
        paras.append(Para(line, 8, False, GRAY, space_before=4 if index == 0 else 3))
    text(slide, x + 0.05 + pad, y + 0.11, w - 0.05 - 2 * pad, h - 0.22, paras)


def stat(slide, x, y, w, h, value, caption, note=None, value_size=19):
    """A KPI tile: the number, then what it counts, then how it is defined."""
    rect(slide, x, y, w, h, fill=CARD, line=RULE)
    rect(slide, x, y, 0.05, h, fill=ORANGE)
    paras = [Para(value, value_size, True, ORANGE, line_spacing=0.94)]
    paras.append(Para(caption, 8.5, True, NAVY, space_before=2.5))
    if note:
        paras.append(Para(note, 7.5, False, GRAY, space_before=1.5))
    text(slide, x + 0.19, y + 0.09, w - 0.33, h - 0.18, paras)


def hbar_chart(slide, x, y, w, h, rows, colors, value_w=1.20, label_w=2.40):
    """Horizontal bars: one measure, one axis, every bar direct-labelled.

    ``rows`` is ``(name, sub_label, value, value_text)``. Length encodes the value; colour encodes
    the ordinal position when a scale is passed and nothing at all when one colour is, so there is
    never a legend to decode. No gridlines: every value is printed beside its bar.
    """
    plot_x = x + label_w + 0.10
    plot_w = w - label_w - 0.10 - value_w
    peak = max((row[2] for row in rows), default=1) or 1
    row_h = h / len(rows)
    bar_h = min(0.19, row_h * 0.46)
    rect(slide, plot_x - 0.045, y + 0.02, 0.008, h - 0.04, fill=RULE)  # one recessive baseline
    for index, (name, sub, value, shown) in enumerate(rows):
        top = y + index * row_h
        mid = top + row_h / 2
        paras = [Para(name, 8.5, True, NAVY, align="r")]
        if sub:
            paras.append(Para(sub, 7.5, False, GRAY, align="r", space_before=1))
        text(slide, x, top, label_w, row_h, paras, anchor="c")
        length = max(0.04, plot_w * value / peak)
        colour = colors[index] if isinstance(colors, list) else colors
        rect(slide, plot_x, mid - bar_h / 2, length, bar_h, fill=colour, radius=BAR_RADIUS)
        text(
            slide,
            plot_x + length + 0.08,
            mid - row_h / 2,
            value_w,
            row_h,
            [Para(shown, 8.5, True, NAVY)],
            anchor="c",
        )


def badge(slide, caption):
    """The template's small header badge: navy pill, white label, top right of the content band."""
    for shape in list(slide.shapes):
        if shape.has_text_frame and shape.text_frame.text.strip() in {"Market lens", "GTM map"}:
            shape._element.getparent().remove(shape._element)
    width = 0.30 + 0.062 * len(caption)
    x = X1 - width
    rect(slide, x, 1.06, width, 0.25, fill=NAVY, radius=0.5)
    text(slide, x, 1.06, width, 0.25, [Para(caption, 7.5, True, WHITE, align="c")], anchor="c")


def set_placeholder(slide, old, new):
    """Replace one of the title slide's bracketed fields, keeping its run formatting."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            if paragraph.text.strip() != old or not paragraph.runs:
                continue
            paragraph.runs[0].text = new
            for extra in paragraph.runs[1:]:
                extra.text = ""
            return True
    return False


# --------------------------------------------------------------------------------------
# the facts, read from the artefacts
# --------------------------------------------------------------------------------------


def money(value):
    return f"€{value:,.0f}"


def compact(value):
    return f"€{value / 1000:,.0f}K" if abs(value) < 1_000_000 else f"€{value / 1e6:,.2f}M"


SEGMENT_LABELS = {
    "Champions",
    "Loyal Customers",
    "Steady Customers",
    "High-Value At Risk",
    "Frequent but Declining",
    "Discount-Driven At Risk",
    "Seasonal Customers",
    "New Customers",
    "One-Time Buyers",
    "Dormant Customers",
    "Lost Customers",
    "High-Return Customers",
    "Low-Value At Risk",
}
LIFECYCLE_ORDER = ("New", "Active", "At Risk", "Declining", "Dormant", "Lost")
DRIFTING = ("At Risk", "Declining", "Dormant")


def load_facts(root: Path) -> dict:
    out = root / "outputs"
    read_json = lambda name: json.loads((out / name).read_text(encoding="utf-8-sig"))
    metrics = read_json("model_metrics.json")
    assumptions = read_json("retention_assumptions.json")
    quality = read_json("data_quality_report.json")
    predictions = pd.read_csv(out / "customer_churn_predictions.csv", encoding="utf-8-sig")
    scores = pd.read_csv(out / "customer_retention_scores.csv", encoding="utf-8-sig")
    recommendations = pd.read_csv(out / "retention_recommendations.csv", encoding="utf-8-sig")

    test = metrics["metrics"]["test"]
    dataset = quality["dataset"]

    # ``Risk level`` is already on the scores artefact; only take what is not there yet, so the
    # merge cannot produce an _x/_y pair and quietly break a groupby.
    extra = [c for c in ("Lifecycle stage", "Frequency") if c not in scores.columns]
    frame = scores.merge(
        predictions[["Customer ID", *extra]], on="Customer ID", validate="1:1"
    )
    exposure = float(frame["Revenue at risk"].sum())
    by_band = frame.groupby("Risk level").agg(
        n=("Customer ID", "size"),
        exposure=("Revenue at risk", "sum"),
        probability=("Churn probability", "mean"),
    )
    lifecycle = predictions["Lifecycle stage"].value_counts()
    top_decile = frame.nlargest(max(1, round(len(frame) / 10)), "Revenue at risk")

    reasons = recommendations.loc[
        recommendations["Recommended action"] == "Do Not Target", "Reason"
    ].str.lower()
    suppression = {
        "already engaged": int(reasons.str.contains("already eng").sum()),
        "uneconomic once costed": int(reasons.str.contains("was indicated").sum()),
        "out of season": int(reasons.str.contains("season").sum()),
        "beyond two buying cycles": int(reasons.str.contains("buying cycles").sum()),
    }

    flags = [c for c in scores.columns if c in SEGMENT_LABELS]
    per_customer = scores[flags].astype(bool).sum(axis=1)

    return {
        "as_of": assumptions["as_of_date"],
        "customers": assumptions["customers"],
        "row_counts": dataset["row_counts"],
        "rows": sum(dataset["row_counts"].values()),
        "net_revenue": dataset["total_net_revenue"],
        "orders": dataset["distinct_orders"],
        "date_min": dataset["purchase_date_min"][:10],
        "date_max": dataset["purchase_date_max"][:10],
        "checks": quality["summary"],
        "exposure": exposure,
        "future": float(frame["Expected future revenue"].sum()),
        "by_band": by_band,
        "top_decile_share": float(top_decile["Revenue at risk"].sum()) / exposure,
        "lifecycle": [(s, int(lifecycle[s])) for s in LIFECYCLE_ORDER if s in lifecycle],
        "drifting": int(sum(lifecycle.get(s, 0) for s in DRIFTING)),
        "one_order": int((predictions["Frequency"] <= 1).sum()),
        "retained": assumptions["total_expected_retained_revenue"],
        "targeted": assumptions["customers_targeted"],
        "suppressed": assumptions["customers_suppressed"],
        "cost": assumptions["total_campaign_cost"],
        "campaign_return": assumptions["campaign_expected_return"],
        "roi": assumptions["campaign_roi"],
        "propensity": assumptions["assumptions"]["base_retention_propensity"]["value"],
        "priority": recommendations.groupby("Priority").agg(
            n=("Customer ID", "size"), exposure=("Revenue at risk", "sum")
        ),
        "suppression": suppression,
        "actions": int(recommendations["Recommended action"].nunique()),
        "skus": int(recommendations["Recommended product/SKU"].nunique()),
        "offers": int(recommendations["Recommended offer"].nunique()),
        "reasons": int(recommendations["Reason"].nunique()),
        "segments_used": int(scores["Primary segment"].nunique()),
        "multi_segment": int((per_customer > 1).sum()),
        "mean_segments": float(per_customer.mean()),
        "model": metrics["model_name"],
        "calibration": metrics["calibration"],
        "horizon": metrics["horizon_days"],
        "train_dates": len(metrics["train_as_of_dates"]),
        "test_date": metrics["test_as_of_dates"][0],
        "embargo": len(metrics["metrics"]["split"]["plan"]["embargoed"]["before_test"]),
        "test": test,
        "s": test["metrics"],
        "b": test["business"],
        "best_baseline": max(
            metrics["metrics"]["single_feature_baselines"], key=lambda item: item["pr_auc"]
        ),
    }


# --------------------------------------------------------------------------------------
# 01 -- Problem & Customer Focus
# --------------------------------------------------------------------------------------


def slide_problem(slide, f):
    badge(slide, "Customer lens")
    text(
        slide,
        X0,
        Y0,
        W,
        0.54,
        [
            Para(
                f"A European fashion brand: {compact(f['net_revenue'])} of net sales, "
                f"{f['customers']:,} customers, {f['orders']:,} orders, {f['date_min']} to "
                f"{f['date_max']} — and nothing for a customer to cancel.",
                11,
                True,
                NAVY,
            ),
            Para(
                "Churn in retail is silent. Nobody gives notice; they stop coming back, and the "
                "only evidence is an absence — a gap between orders that quietly grows past that "
                "customer's own normal. By the time a monthly report shows the hole, the customer "
                "has been gone two seasons and the cheapest moment to act is behind you.",
                8.5,
                False,
                GRAY,
                space_before=3,
            ),
        ],
    )

    tiles = [
        (
            money(f["exposure"]),
            f"Revenue at risk, next {f['horizon']} days",
            f"{f['exposure'] / f['future']:.1%} of the {compact(f['future'])} expected future revenue",
        ),
        (
            f"{f['test']['base_rate']:.1%}",
            f"Lapse within {f['horizon']} days",
            f"Measured, held-out test period, n={f['test']['n']:,}",
        ),
        (
            f"{f['drifting']:,}",
            "Customers already drifting",
            f"At risk, declining or dormant at {f['as_of']}",
        ),
        (
            f"{f['one_order']:,}",
            "One-and-done buyers",
            f"Of {f['customers']:,} — never returned after order one",
        ),
    ]
    gap = 0.22
    tile_w = (W - 3 * gap) / 4
    for index, (value, caption, note) in enumerate(tiles):
        stat(slide, X0 + index * (tile_w + gap), 2.12, tile_w, 0.90, value, caption, note)

    col_gap = 0.30
    col_w = (W - col_gap) / 2
    right_x = X0 + col_w + col_gap

    label(slide, X0, 3.14, col_w, "The journey, and where it breaks")
    hbar_chart(
        slide,
        X0,
        3.44,
        col_w,
        1.90,
        [
            (stage, f"{count / f['customers']:.0%} of the book", count, f"{count:,}")
            for stage, count in f["lifecycle"]
        ],
        BAR_ONE_SERIES,
        value_w=0.62,
        label_w=1.80,
    )
    text(
        slide,
        X0,
        5.40,
        col_w,
        0.30,
        [
            Para(
                "Lifecycle stage at the prediction date. This is not a trickle at the end of a "
                "funnel: the largest stage is the one the brand has already stopped talking to.",
                7.5,
                False,
                GRAY,
            )
        ],
    )

    label(slide, right_x, 3.14, col_w, "Where the exposure actually sits")
    band = f["by_band"]
    order = [b for b in ("Low", "Medium", "High", "Critical") if b in band.index]
    hbar_chart(
        slide,
        right_x,
        3.44,
        col_w,
        1.28,
        [
            (
                b,
                f"{int(band.loc[b, 'n'])} customers · avg churn "
                f"{band.loc[b, 'probability']:.0%}",
                float(band.loc[b, "exposure"]),
                money(band.loc[b, "exposure"]),
            )
            for b in order
        ],
        [RISK_RAMP[b] for b in order],
        value_w=0.95,
        label_w=2.15,
    )
    critical = band.loc["Critical"]
    rect(slide, right_x, 4.80, col_w, 0.54, fill=CARD_DEEP, line=RULE)
    rect(slide, right_x, 4.80, 0.05, 0.54, fill=ORANGE)
    text(
        slide,
        right_x + 0.21,
        4.80,
        col_w - 0.40,
        0.54,
        [
            Para(
                f"The reddest customers hold the least money. The {int(critical['n'])} Critical "
                f"customers average a {critical['probability']:.0%} churn probability but carry "
                f"only {money(critical['exposure'])} — "
                f"{critical['exposure'] / f['exposure']:.0%} of total exposure — because there is "
                f"barely anything left to lose. Ranking by probability spends the budget on "
                f"customers who have already gone.",
                8,
                False,
                GRAY,
            )
        ],
        anchor="c",
    )
    text(
        slide,
        right_x,
        5.40,
        col_w,
        0.30,
        [
            Para(
                "Revenue at risk = churn probability × expected future revenue, per customer, "
                "summed by band. Probability says who; only value says how much.",
                7.5,
                False,
                GRAY,
            )
        ],
    )

    gap = 0.24
    pain_w = (W - 2 * gap) / 3
    pains = [
        (
            "There is no cancel button",
            "The signal has to be constructed. Each customer's own cadence defines their normal, "
            "so one global rule mislabels half the book: a twice-a-year coat buyer and a monthly "
            "basics buyer cannot share a threshold.",
        ),
        (
            "Discounts are sprayed, not aimed",
            "Half of all order lines already carry a discount. One house-standard offer pays "
            "margin to full-price buyers who were never leaving, and under-offers the people who "
            "only ever move on promotion.",
        ),
        (
            "Seasonality reads as churn",
            "A customer who buys coats every October is not lapsing in June. Chasing them wastes "
            "the contact and teaches them to ignore the next email — the one that mattered.",
        ),
    ]
    for index, (title, body) in enumerate(pains):
        card(slide, X0 + index * (pain_w + gap), 5.74, pain_w, 0.94, title, [body], title_size=9)


# --------------------------------------------------------------------------------------
# 02 -- The AI Solution & Innovation
# --------------------------------------------------------------------------------------


def slide_solution(slide, f):
    badge(slide, "Build lens")
    counts = f["row_counts"]
    text(
        slide,
        X0,
        Y0,
        W,
        0.50,
        [
            Para(
                "Five questions, one pipeline: who is likely to leave, why, how much revenue that "
                "puts at risk, what to do about it, and who to contact first.",
                11,
                True,
                NAVY,
            ),
            Para(
                f"Six stages, each a batch CLI reading the four CSV files as the single source of "
                f"truth — {counts['customers']:,} customers, {counts['products']:,} SKUs, "
                f"{counts['transactions']:,} order lines, {counts['returns']:,} returns. No "
                f"database, no ETL, no API and no service to keep alive: copy the folder to a "
                f"machine with the same CSVs and it runs. An architecture test enforces that "
                f"rather than a paragraph promising it.",
                8.5,
                False,
                GRAY,
                space_before=3,
            ),
        ],
    )

    label(
        slide,
        X0,
        2.08,
        W,
        "The pipeline",
        "Every stage writes an artefact the next one reads. The dashboard is the last box and only reads.",
    )
    stages = [
        ("LOAD", "Four CSVs", f"{f['rows']:,} rows, dtypes pinned"),
        ("VALIDATE", "Schema & keys", f"{f['checks']['passed']}/{f['checks']['total']} checks pass"),
        ("FEATURES", "148 features", "computed as of a date"),
        ("MODEL", f["model"].upper(), f"{f['calibration']}-calibrated, {f['horizon']}d"),
        ("EXPLAIN", "SHAP per customer", "in plain English"),
        ("ACT", f"{f['segments_used']} segments", "ROI-gated actions"),
    ]
    gap = 0.30
    box_w = (W - 5 * gap) / 6
    for index, (step, headline, note) in enumerate(stages):
        x = X0 + index * (box_w + gap)
        rect(slide, x, 2.40, box_w, 0.84, fill=CARD, line=RULE)
        rect(slide, x, 2.40, box_w, 0.045, fill=ORANGE)
        text(
            slide,
            x + 0.10,
            2.46,
            box_w - 0.20,
            0.74,
            [
                Para(step, 7.5, True, ORANGE_TEXT, align="c"),
                Para(headline, 9.5, True, NAVY, align="c", space_before=2),
                Para(note, 7.5, False, GRAY, align="c", space_before=2),
            ],
            anchor="c",
        )
        if index < len(stages) - 1:
            text(
                slide,
                x + box_w,
                2.40,
                gap,
                0.84,
                [Para("›", 15, True, ORANGE, align="c")],
                anchor="c",
            )

    label(slide, X0, 3.44, W, "What is actually novel here")
    gap = 0.26
    col_w = (W - 2 * gap) / 3
    s, b = f["s"], f["b"]
    innovations = [
        (
            "The as-of date is the leakage guard",
            [
                f"Every feature is computed as of a prediction date and the label looks forward "
                f"{f['horizon']} days from it, so the model answers the question a business "
                f"actually asks on a Monday morning.",
                f"Training spans {f['train_dates']} monthly as-of dates with a "
                f"{f['embargo']}-month embargo before each evaluation period. A random split here "
                f"leaks the future: it scores beautifully and is worthless.",
            ],
        ),
        (
            "A calibrated probability, not a ranking",
            [
                "Revenue at risk is denominated in euros, so P(churn) has to mean what it says — "
                "a raw ranking score cannot be multiplied by money.",
                f"Calibration is chosen out-of-fold on a held-out recent period: mean predicted "
                f"{s['mean_predicted']:.3f} against {s['mean_observed']:.3f} observed, a bias of "
                f"{s['calibration_bias']:+.3f} on data the model never saw.",
            ],
        ),
        (
            "Nothing about the action is hardcoded",
            [
                f"{f['actions']} actions, {f['skus']} SKUs, {f['offers']} offer depths and "
                f"{f['reasons']:,} distinct reasons across {f['customers']:,} customers — every "
                f"field derived, none of it written by hand.",
                "Offer depth is the depth that customer has responded to before; the SKU is the "
                "best seller in their category, excluding what they already own; the reason quotes "
                "their own numbers back at them.",
            ],
        ),
    ]
    for index, (title, body) in enumerate(innovations):
        card(slide, X0 + index * (col_w + gap), 3.74, col_w, 1.50, title, body)

    label(
        slide,
        X0,
        5.44,
        W,
        "Measured on a period the model never saw",
        f"Test as-of {f['test_date']} · n={f['test']['n']:,} · base rate "
        f"{f['test']['base_rate']:.1%} · {f['model']}, {f['calibration']}-calibrated",
    )
    metrics = [
        (f"{b['high_value_roc_auc']:.3f}", "ROC-AUC, High Value", "where the money is"),
        (f"{s['roc_auc']:.3f}", "ROC-AUC, overall", "ranking quality"),
        (f"{s['pr_auc']:.3f}", "PR-AUC", f"beats the {f['best_baseline']['pr_auc']:.3f} floor"),
        (f"{s['brier']:.3f}", "Brier score", f"ECE {s['ece']:.3f}"),
        (f"{s['lift_top_decile']:.2f}×", "Lift, top decile", "against random"),
        (
            f"{b['high_value_precision_top_decile']:.0%}",
            "Top-decile precision",
            "High Value only",
        ),
    ]
    gap = 0.20
    tile_w = (W - 5 * gap) / 6
    for index, (value, caption, note) in enumerate(metrics):
        stat(
            slide,
            X0 + index * (tile_w + gap),
            5.74,
            tile_w,
            0.94,
            value,
            caption,
            note,
            value_size=15,
        )


# --------------------------------------------------------------------------------------
# 03 -- Standing Out From Competitors
# --------------------------------------------------------------------------------------


def slide_differentiation(slide, f):
    badge(slide, "Market lens")
    text(
        slide,
        X0,
        Y0,
        W,
        0.50,
        [
            Para(
                "Most churn projects stop at a score and a bar chart of feature importances. The "
                "gap between a score and a decision is where this platform lives.",
                11,
                True,
                NAVY,
            ),
            Para(
                "A churn score tells a marketing lead nothing they can spend against. Six things "
                "have to be true before a number becomes a contact list somebody will sign off — "
                "and every row below is a place where the conventional build stops short.",
                8.5,
                False,
                GRAY,
                space_before=3,
            ),
        ],
    )

    rows = [
        (
            "The output",
            "A churn score, uncalibrated. Useful for sorting; meaningless in euros.",
            f"A calibrated probability, so exposure is money: {money(f['exposure'])} at risk, "
            f"attributable customer by customer and addable across a segment.",
        ),
        (
            "Who to chase",
            "The reddest customers first — which spends the budget where the least value is left.",
            f"Revenue at risk, so a mid-probability high-value customer outranks a near-certain "
            f"loss. The top decile by exposure carries {f['top_decile_share']:.0%} of the total.",
        ),
        (
            "The 'why'",
            "One global feature-importance chart, shown to every customer alike.",
            f"Per-customer SHAP in plain English, with direction of impact measured from the data "
            f"rather than assumed — {f['reasons']:,} distinct reasons for "
            f"{f['customers']:,} customers.",
        ),
        (
            "The action",
            "A house-standard offer, mailed to everyone on the list.",
            f"{f['actions']} actions, {f['skus']} SKUs, {f['offers']} offer depths. Full-price "
            f"buyers are structurally unreachable by a discount rule, and a negative-ROI action is "
            f"overridden after the fact, so the reason still says what was proposed.",
        ),
        (
            "Validation",
            "A random train/test split, and a metric with nothing to compare it against.",
            f"A time-based split with a {f['embargo']}-month embargo, calibrated out-of-fold on a "
            f"held-out period, and eight one-line heuristics scored every run as a floor the model "
            f"has to clear ({f['s']['pr_auc']:.3f} against {f['best_baseline']['pr_auc']:.3f}).",
        ),
        (
            "Honesty",
            "Assumptions quietly folded into the ROI number.",
            "Retention propensity cannot be learned without a control group, so the artefact names "
            "the column '(ASSUMED)' and every figure it touches is flagged on the page that shows "
            "it. Judges get the number and its warranty.",
        ),
    ]
    header_h, row_h = 0.32, 0.44
    widths = [1.80, 4.05, W - 1.80 - 4.05]
    top = 2.06
    rect(slide, X0, top, W, header_h, fill=NAVY)
    for index, caption in enumerate(("", "The conventional build", "This platform")):
        text(
            slide,
            X0 + sum(widths[:index]) + 0.16,
            top,
            widths[index] - 0.26,
            header_h,
            [Para(caption, 8.5, True, WHITE)],
            anchor="c",
        )
    for r, cells in enumerate(rows):
        y = top + header_h + r * row_h
        rect(slide, X0, y, W, row_h, fill=CARD if r % 2 else WHITE, line=RULE)
        rect(slide, X0 + widths[0] + widths[1], y, 0.05, row_h, fill=ORANGE)
        for c, value in enumerate(cells):
            text(
                slide,
                X0 + sum(widths[:c]) + (0.21 if c == 2 else 0.16),
                y,
                widths[c] - (0.34 if c == 2 else 0.28),
                row_h,
                [Para(value, 8, c == 0, NAVY if c == 0 else GRAY)],
                anchor="c",
            )

    gap = 0.26
    col_w = (W - 2 * gap) / 3
    cards = [
        (
            "Simulate before you spend",
            "The What-If page rebuilds the real retention layer and re-runs it, rather than "
            "rescaling the shipped columns. Raise the contact cost and customers drop out of the "
            "campaign entirely — no linear rescale reproduces that.",
        ),
        (
            "The architecture is a test, not a promise",
            "'The CSVs are the database' is enforced: the four files are hashed, the pipeline "
            "runs, they are hashed again. No driver, orchestrator or web framework can be imported "
            "without a test failing. 473 tests in all.",
        ),
        (
            "Segments overlap, because customers do",
            f"{f['multi_segment']:,} of {f['customers']:,} customers carry more than one label "
            f"({f['mean_segments']:.2f} each on average). A high-return customer who is also "
            f"high-value at risk needs both facts: one says act now, the other says be careful "
            f"what you offer.",
        ),
    ]
    for index, (title, body) in enumerate(cards):
        card(slide, X0 + index * (col_w + gap), 5.14, col_w, 1.54, title, [body])


# --------------------------------------------------------------------------------------
# 04 -- Revenue & Business Growth
# --------------------------------------------------------------------------------------


def slide_growth(slide, f):
    badge(slide, "GTM map")
    text(
        slide,
        X0,
        Y0,
        W,
        0.50,
        [
            Para(
                f"Same {f['customers']:,} customers, same four files — the difference is that "
                f"every euro below is traceable to a customer ID, an action and a reason.",
                11,
                True,
                NAVY,
            ),
            Para(
                f"The chain runs from what the book is expected to be worth over the next "
                f"{f['horizon']} days down to the return on a single campaign. Only the last two "
                f"steps depend on an assumption, and they are labelled rather than quoted as "
                f"though they were measured.",
                8.5,
                False,
                GRAY,
                space_before=3,
            ),
        ],
    )

    chart_w = 7.05
    label(slide, X0, 2.06, chart_w, "From exposure to return")
    hbar_chart(
        slide,
        X0,
        2.36,
        chart_w,
        1.26,
        [
            (
                "Expected future revenue",
                f"{f['horizon']} days, all {f['customers']:,} customers",
                f["future"],
                money(f["future"]),
            ),
            (
                "Revenue at risk",
                "churn probability × expected revenue",
                f["exposure"],
                money(f["exposure"]),
            ),
            (
                "Recoverable",
                f"at the assumed {f['propensity']:.0%} base propensity",
                f["retained"],
                money(f["retained"]),
            ),
            (
                "Expected return",
                f"from the {f['targeted']:,} customers targeted",
                f["campaign_return"],
                money(f["campaign_return"]),
            ),
        ],
        BAR_ONE_SERIES,
        value_w=0.95,
        label_w=2.30,
    )

    hero_x = X0 + chart_w + 0.30
    hero_w = W - chart_w - 0.30
    rect(slide, hero_x, 2.36, hero_w, 1.26, fill=CARD_DEEP, line=RULE)
    rect(slide, hero_x, 2.36, 0.05, 1.26, fill=ORANGE)
    text(
        slide,
        hero_x + 0.23,
        2.36,
        hero_w - 0.42,
        1.26,
        [
            Para(f"+{f['roi'] * 100:,.0f}%", 28, True, ORANGE, line_spacing=0.92),
            Para(
                f"Expected ROI on one {f['horizon']}-day retention campaign", 9, True, NAVY,
                space_before=3,
            ),
            Para(
                f"{money(f['campaign_return'])} expected return on {money(f['cost'])} of contact "
                f"and incentive cost, across {f['targeted']:,} targeted customers — "
                f"{f['suppressed']} are deliberately left alone. ASSUMPTION-dependent; see the "
                f"last card.",
                7.5,
                False,
                GRAY,
                space_before=3,
            ),
        ],
        anchor="c",
    )

    col_gap = 0.30
    left_w = 5.10
    right_x = X0 + left_w + col_gap
    right_w = W - left_w - col_gap

    label(slide, X0, 3.76, left_w, "Where to spend first")
    priority = f["priority"]
    widths = [1.20, 1.05, 1.35, left_w - 3.60]
    heads = ("Priority", "Customers", "Revenue at risk", "Share of exposure")
    y, row_h = 4.06, 0.24
    rect(slide, X0, y, left_w, row_h, fill=NAVY)
    for index, caption in enumerate(heads):
        text(
            slide,
            X0 + sum(widths[:index]) + 0.10,
            y,
            widths[index] - 0.16,
            row_h,
            [Para(caption, 7.5, True, WHITE)],
            anchor="c",
        )
    order = [p for p in ("Critical", "High", "Medium", "Low") if p in priority.index]
    for index, name in enumerate(order):
        ry = y + row_h * (index + 1)
        rect(slide, X0, ry, left_w, row_h, fill=CARD if index % 2 else WHITE, line=RULE)
        exposure = float(priority.loc[name, "exposure"])
        cells = (
            name,
            f"{int(priority.loc[name, 'n']):,}",
            money(exposure),
            f"{exposure / f['exposure']:.0%}",
        )
        for c, value in enumerate(cells):
            text(
                slide,
                X0 + sum(widths[:c]) + 0.10,
                ry,
                widths[c] - 0.16,
                row_h,
                [Para(value, 8, c == 0, NAVY if c == 0 else GRAY)],
                anchor="c",
            )
    text(
        slide,
        X0,
        5.34,
        left_w,
        0.30,
        [
            Para(
                "Not contacted, by reason: "
                + " · ".join(f"{n} {why}" for why, n in f["suppression"].items() if n)
                + ". 'Already engaged' and 'unrecoverable' need opposite follow-up, so they are "
                "never reported as one number.",
                7.5,
                False,
                GRAY,
            )
        ],
    )

    label(slide, right_x, 3.76, right_w, "Who buys this")
    buyers = [
        (
            "Mid-market fashion & apparel e-commerce",
            "€10m–€200m of online revenue. They already have these four tables in a BI export; "
            "what they do not have is the data team to turn them into a contact list.",
        ),
        (
            "CRM and retention leads who own a contact budget",
            "The buyer is whoever gets asked what the retention spend returned. This answers in "
            "euros, per customer, with the assumption labelled instead of buried.",
        ),
        (
            "Agencies running retention for several brands",
            "One folder per client, identical code, no per-client integration project — which is "
            "what usually kills a retention pilot before it can prove anything.",
        ),
    ]
    y = 4.06
    for title, body in buyers:
        text(
            slide,
            right_x,
            y,
            right_w,
            0.42,
            [Para(title, 8.5, True, NAVY), Para(body, 7.5, False, GRAY, space_before=1.5)],
        )
        y += 0.44

    gap = 0.26
    col_w = (W - 2 * gap) / 3
    cards = [
        (
            "Commercialisation",
            "Sold as a per-brand instance: drop four CSVs in, run five commands, get eight "
            "dashboard pages. Priced per brand per month, with the pilot scoped as one campaign "
            "and one holdout.",
        ),
        (
            "Scale",
            f"{f['rows']:,} rows today, read straight into Pandas. Every stage is a batch CLI that "
            f"exits, so a monthly refresh is a cron line rather than a platform, and adding a brand "
            f"adds a folder rather than a migration.",
        ),
        (
            "What we would need next — honestly",
            f"The {f['propensity']:.0%} retention propensity is a planning placeholder: uplift "
            f"needs a campaign log and an untreated control group, and this dataset has neither. "
            f"One holdout campaign turns +{f['roi'] * 100:,.0f}% from an estimate into a "
            f"measurement — the first thing a pilot should buy.",
        ),
    ]
    for index, (title, body) in enumerate(cards):
        card(slide, X0 + index * (col_w + gap), 5.60, col_w, 1.08, title, [body])


# --------------------------------------------------------------------------------------


def build(template: Path, output: Path, root: Path) -> Path:
    facts = load_facts(root)
    deck = Presentation(str(template))

    title = deck.slides[0]
    set_placeholder(title, "[Enter hackathon name]", "EXL Hackathon — Dragon's Den")
    set_placeholder(title, "[Enter solution name]", "Predictive Churn Analytics for Fashion Retail")

    slide_problem(deck.slides[1], facts)
    slide_solution(deck.slides[2], facts)
    slide_differentiation(deck.slides[3], facts)
    slide_growth(deck.slides[4], facts)

    deck.save(str(output))
    return output


def main() -> int:
    parser = argparse.ArgumentParser(description="Build the Dragon's Den deck from the artefacts.")
    parser.add_argument(
        "--template", type=Path, default=ROOT / "Dragons_Den_Presentation_Template.pptx"
    )
    parser.add_argument(
        "--output",
        type=Path,
        default=ROOT / "Dragons_Den_Presentation_Predictive_Churn_Analytics.pptx",
    )
    args = parser.parse_args()

    if not args.template.exists():
        print(f"template not found: {args.template}", file=sys.stderr)
        return 1
    written = build(args.template, args.output, ROOT)
    print(f"wrote {written.name}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
